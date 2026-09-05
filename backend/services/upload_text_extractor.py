"""Extract bounded text from contribution uploads before Hermes sees them."""
from __future__ import annotations

import io
from pathlib import Path

MAX_EXTRACTED_CHARACTERS = 200_000


def _bounded(text: str) -> str:
    value = text.strip()
    if not value:
        raise ValueError("uploaded file contains no extractable text")
    if len(value) > MAX_EXTRACTED_CHARACTERS:
        raise ValueError("uploaded file extracted text exceeds contribution limit")
    return value


def extract_uploaded_text(data: bytes, *, filename: str, content_type: str) -> str:
    suffix = Path(filename).suffix.lower()
    mime = content_type.split(";", 1)[0].strip().lower()
    if mime.startswith("text/") or suffix in {".md", ".txt", ".csv", ".json", ".yaml", ".yml", ".html"}:
        return _bounded(data.decode("utf-8", errors="replace"))
    if suffix == ".pdf" or mime == "application/pdf":
        from pypdf import PdfReader
        return _bounded("\n\n".join(page.extract_text() or "" for page in PdfReader(io.BytesIO(data)).pages))
    if suffix == ".docx" or mime == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
        from docx import Document
        document = Document(io.BytesIO(data))
        parts = [p.text for p in document.paragraphs]
        parts.extend("\t".join(cell.text for cell in row.cells)
                     for table in document.tables for row in table.rows)
        return _bounded("\n".join(parts))
    if suffix == ".xlsx" or mime == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":
        from openpyxl import load_workbook
        workbook = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
        lines = []
        for sheet in workbook.worksheets:
            lines.append(f"# {sheet.title}")
            lines.extend("\t".join("" if value is None else str(value) for value in row)
                         for row in sheet.iter_rows(values_only=True))
        return _bounded("\n".join(lines))
    if suffix == ".pptx" or mime == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
        from pptx import Presentation
        presentation = Presentation(io.BytesIO(data))
        return _bounded("\n".join(
            shape.text for slide in presentation.slides for shape in slide.shapes
            if hasattr(shape, "text") and shape.text
        ))
    raise ValueError("uploaded file type has no safe text extractor")
